import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from execution.render_pptx import render_presentation
from tests.support import make_metadata, make_valid_payload


class RenderPresentationTests(unittest.TestCase):
    def test_renders_editable_deck_with_expected_slide_titles(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = Path(temp_dir) / "meeting_deck.pptx"

            render_presentation(make_valid_payload(), make_metadata(), output_path)

            presentation = Presentation(output_path)
            titles = [slide.shapes.title.text for slide in presentation.slides]

            self.assertEqual(5, len(presentation.slides))
            self.assertEqual(
                [
                    "BrightLane Retail AI Automation Discovery Call",
                    "Executive Summary",
                    "Objectives",
                    "Action Items",
                    "Next Steps and Risks",
                ],
                titles,
            )


if __name__ == "__main__":
    unittest.main()

