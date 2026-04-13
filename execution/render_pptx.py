from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation


def _add_bullets(text_frame, lines: list[str]) -> None:
    text_frame.clear()
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0


def render_presentation(payload: dict[str, Any], metadata: dict[str, Any], output_path: str | Path) -> Path:
    presentation = Presentation()

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = metadata["meeting_title"]
    title_slide.placeholders[1].text = (
        f"{metadata['meeting_date']}\nReview required: {metadata.get('review_required', True)}"
    )

    summary_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    summary_slide.shapes.title.text = "Executive Summary"
    _add_bullets(
        summary_slide.placeholders[1].text_frame,
        [
            payload["executive_summary"],
            f"Participants: {', '.join(metadata.get('participants', []))}",
        ],
    )

    objectives_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    objectives_slide.shapes.title.text = "Objectives"
    _add_bullets(
        objectives_slide.placeholders[1].text_frame,
        [f"{item['title']}: {item['rationale']}" for item in payload["objectives"]],
    )

    actions_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    actions_slide.shapes.title.text = "Action Items"
    _add_bullets(
        actions_slide.placeholders[1].text_frame,
        [
            f"{item['title']} | Owner: {item['owner']} | Due: {item['due_window']}"
            for item in payload["action_items"]
        ],
    )

    next_steps_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    next_steps_slide.shapes.title.text = "Next Steps and Risks"
    _add_bullets(
        next_steps_slide.placeholders[1].text_frame,
        ["Next steps:"] + payload["next_steps"] + ["Risks:"] + payload["risks"],
    )

    output_file = Path(output_path)
    output_file.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_file)
    return output_file

